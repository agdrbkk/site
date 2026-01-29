from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation_v2():
    prs = Presentation()

    # Define Theme Colors
    BG_DARK = RGBColor(18, 18, 18)      # Almost Black
    ACCENT_RED = RGBColor(204, 0, 0)    # Crimson Red
    ACCENT_GOLD = RGBColor(197, 160, 90) # Muted Gold
    TEXT_WHITE = RGBColor(240, 240, 240)
    TEXT_GREY = RGBColor(180, 180, 180)

    def format_text(paragraph, text, size, color, bold=False, align=PP_ALIGN.LEFT):
        paragraph.text = text
        paragraph.font.size = size
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        paragraph.alignment = align
        paragraph.font.name = 'Arial'

    def add_design_header(slide, title_text):
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

        # Red decorative strip on the left
        left_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
        left_strip.fill.solid()
        left_strip.fill.fore_color.rgb = ACCENT_RED
        left_strip.line.fill.background() # No line

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(1))
        format_text(title_box.text_frame.paragraphs[0], title_text.upper(), Pt(36), TEXT_WHITE, True)

        # Gold Underline
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(3), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_GOLD
        line.line.fill.background()

    def add_footer(slide):
        # Footer text
        footer = slide.shapes.add_textbox(Inches(0.8), Inches(7), Inches(5), Inches(0.5))
        format_text(footer.text_frame.paragraphs[0], "SHABU-MALA INTERIOR PROPOSAL | 2026", Pt(10), TEXT_GREY)

    # ==========================
    # SLIDE 1: TITLE SLIDE
    # ==========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Dark Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_DARK

    # Big Red Block graphics
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_RED
    shape.line.fill.background()

    # Main Title (centered in the red block or below)
    # Let's put text below the block for contrast or inside. 
    # Let's go for a layout where the top third is red.
    
    # Title Text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.5))
    p = txBox.text_frame.paragraphs[0]
    format_text(p, "INTERIOR DESIGN\nPROPOSAL", Pt(54), TEXT_WHITE, True, PP_ALIGN.LEFT)
    
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    p2 = txBox2.text_frame.paragraphs[0]
    format_text(p2, "CHINESE SHABU-MALA RESTAURANT", Pt(28), ACCENT_GOLD, True, PP_ALIGN.LEFT)

    # Details
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    p3 = txBox3.text_frame.paragraphs[0]
    format_text(p3, "THEME: RED & BLACK AESTHETIC\nPREPARED FOR: CLIENT", Pt(18), TEXT_GREY, False, PP_ALIGN.LEFT)

    # ==========================
    # SLIDE 2: OVERVIEW
    # ==========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_design_header(slide, "Design Concepts")
    add_footer(slide)

    concepts = [
        ("MODERN DARK LUXURY", "Matte Black & Crimson | Mysterious & Premium"),
        ("NEO-CYBER CHINATOWN", "Industrial & Neon | Energetic & Street"),
        ("THE CRIMSON THEATRE", "Wood & Grandeur | Dramatic & Cultural")
    ]

    top_pos = 2.0
    for title, desc in concepts:
        # Number Box
        num_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top_pos), Inches(0.4), Inches(0.4))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = ACCENT_RED
        num_box.line.fill.background()
        
        # Text
        tx = slide.shapes.add_textbox(Inches(1.4), Inches(top_pos-0.15), Inches(8), Inches(0.8))
        p = tx.text_frame.paragraphs[0]
        format_text(p, title, Pt(20), TEXT_WHITE, True)
        
        p_desc = tx.text_frame.add_paragraph()
        format_text(p_desc, desc, Pt(16), TEXT_GREY)
        
        top_pos += 1.3

    # ==========================
    # CONCEPT SLIDES FUNCTION
    # ==========================
    def create_concept_slide(title, mood, materials, refs):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_design_header(s, title)
        add_footer(s)

        # Left Column: Content
        # Mood
        tx_mood = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4), Inches(0.5))
        p = tx_mood.text_frame.paragraphs[0]
        format_text(p, "MOOD & TONE", Pt(14), ACCENT_GOLD, True)
        
        tx_mood_desc = s.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(4), Inches(1))
        p = tx_mood_desc.text_frame.paragraphs[0]
        format_text(p, mood, Pt(16), TEXT_WHITE)

        # Materials
        tx_mat = s.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(4), Inches(0.5))
        p = tx_mat.text_frame.paragraphs[0]
        format_text(p, "KEY MATERIALS", Pt(14), ACCENT_GOLD, True)
        
        tx_mat_desc = s.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(4), Inches(2))
        p = tx_mat_desc.text_frame.paragraphs[0]
        format_text(p, materials, Pt(16), TEXT_WHITE)

        # Right Column: Image Placeholders
        # Frame 1
        frame1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.8), Inches(4), Inches(2.5))
        frame1.fill.solid()
        frame1.fill.fore_color.rgb = RGBColor(40, 40, 40) # Dark grey placeholder
        frame1.line.color.rgb = ACCENT_GOLD
        frame1.line.width = Pt(1)
        frame1.text_frame.text = "PASTE IMAGE HERE\n(Perspective View)"
        frame1.text_frame.paragraphs[0].font.color.rgb = TEXT_GREY

        # Frame 2
        frame2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.5), Inches(4), Inches(2.2))
        frame2.fill.solid()
        frame2.fill.fore_color.rgb = RGBColor(40, 40, 40)
        frame2.line.color.rgb = ACCENT_GOLD
        frame2.line.width = Pt(1)
        frame2.text_frame.text = f"PASTE IMAGE HERE\n(Details/Materials)\n\nRefs: {refs}"
        frame2.text_frame.paragraphs[0].font.size = Pt(10)
        frame2.text_frame.paragraphs[0].font.color.rgb = TEXT_GREY

    # Create the 3 Concept Slides
    create_concept_slide(
        "01: Modern Dark Luxury",
        "Mysterious, Sophisticated, Private Lounge Vibe",
        "- Matte Black Walls\n- Dark Crimson Velvet\n- Hidden Warm Lighting\n- Black Marble",
        "YKC II / Da Ming Ding Ding"
    )

    create_concept_slide(
        "02: Neo-Cyber Chinatown",
        "Energetic, Street Style, Futuristic, Photogenic",
        "- Industrial Concrete\n- Red Neon Signs\n- Metal Mesh\n- Stainless Steel",
        "Tiago Select / Cyberpunk Art"
    )

    create_concept_slide(
        "03: The Crimson Theatre",
        "Dramatic, Cultural, Grand, Cinematic",
        "- Black Wood Structure\n- Giant Red Pillars\n- Fabric Lanterns\n- Symmetrical Layout",
        "Jiu Li Hot Pot"
    )

    # ==========================
    # FINAL SLIDE
    # ==========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_DARK

    # Center box
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.5), Inches(6), Inches(2.5))
    box.fill.background() # No fill
    box.line.color.rgb = ACCENT_GOLD
    box.line.width = Pt(2)

    tx = slide.shapes.add_textbox(Inches(2), Inches(3.2), Inches(6), Inches(1))
    p = tx.text_frame.paragraphs[0]
    format_text(p, "THANK YOU", Pt(48), TEXT_WHITE, True, PP_ALIGN.CENTER)
    
    tx2 = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(1))
    p2 = tx2.text_frame.paragraphs[0]
    format_text(p2, "Ready to bring this vision to life?", Pt(18), TEXT_GREY, False, PP_ALIGN.CENTER)

    prs.save('Mala_Hotpot_Proposal_V2.pptx')
    print("Presentation V2 saved successfully.")

if __name__ == "__main__":
    create_presentation_v2()
