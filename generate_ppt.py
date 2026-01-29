from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Define colors
    BLACK = RGBColor(0, 0, 0)
    DARK_GREY = RGBColor(30, 30, 30)
    RED = RGBColor(200, 20, 30)
    WHITE = RGBColor(255, 255, 255)
    GOLD = RGBColor(212, 175, 55)

    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text(slide, text, left, top, width, height, font_size, color=WHITE, bold=False):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        return txBox

    # 1. Cover Slide
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, BLACK)
    
    # Title
    add_text(slide, "INTERIOR DESIGN PROPOSAL", Inches(1), Inches(2.5), Inches(8), Inches(1), Pt(44), RED, True)
    add_text(slide, "Modern Chinese Shabu-Mala Restaurant", Inches(1), Inches(3.5), Inches(8), Inches(1), Pt(32), WHITE)
    add_text(slide, "Theme: Red & Black Aesthetic", Inches(1), Inches(4.5), Inches(8), Inches(1), Pt(24), GOLD)

    # 2. Concept Overview
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, DARK_GREY)
    add_text(slide, "3 DESIGN CONCEPTS", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(36), RED, True)
    
    concepts = [
        "1. Modern Dark Luxury: เรียบหรู ลึกลับ (Matte Black & Dark Crimson)",
        "2. Neo-Cyber Chinatown: สนุกจัดจ้าน (Industrial Black & Neon Red)",
        "3. The Crimson Theatre: โรงเตี๊ยมร่วมสมัย (Black Wood & Chinese Red)"
    ]
    
    top = 2.0
    for concept in concepts:
        add_text(slide, concept, Inches(1), Inches(top), Inches(8), Inches(1), Pt(24), WHITE)
        top += 1.2

    # 3. Concept 1
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, BLACK)
    add_text(slide, "CONCEPT 01: MODERN DARK LUXURY", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(32), RED, True)
    
    details = "Mood: Mysterious, Sophisticated, High-End\n\n" \
              "Key Elements:\n" \
              "- Matte Black Walls (ผนังดำด้าน)\n" \
              "- Dark Crimson Velvet (กำมะหยี่แดงเลือดหมู)\n" \
              "- Hidden Lighting (ไฟซ่อน/สปอตไลท์)\n\n" \
              "References:\n" \
              "- ArchDaily: YKC II Restaurant\n" \
              "- WorldArchitecture: Da Ming Ding Ding"
    
    tx = add_text(slide, details, Inches(0.5), Inches(1.5), Inches(5), Inches(5), Pt(18), WHITE)
    tx.text_frame.word_wrap = True
    
    # Placeholder for images
    slide.shapes.add_shape(1, Inches(6), Inches(1.5), Inches(3.5), Inches(2.5)).text = "Image Placeholder 1\n(Luxury Dark Interior)"
    slide.shapes.add_shape(1, Inches(6), Inches(4.5), Inches(3.5), Inches(2.5)).text = "Image Placeholder 2\n(Red Velvet & Marble)"

    # 4. Concept 2
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, BLACK)
    add_text(slide, "CONCEPT 02: NEO-CYBER CHINATOWN", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(32), RED, True)
    
    details = "Mood: Energetic, Street Vibe, Photogenic\n\n" \
              "Key Elements:\n" \
              "- Industrial Loft (ปูนเปลือย/โครงเหล็ก)\n" \
              "- Neon Lights (ไฟนีออนดัดตัวอักษรจีน)\n" \
              "- Metal Mesh (ตะแกรงเหล็กฉีก)\n\n" \
              "References:\n" \
              "- Designboom: Tiago Select\n" \
              "- ArtStation: Cyberpunk Restaurant"
    
    tx = add_text(slide, details, Inches(0.5), Inches(1.5), Inches(5), Inches(5), Pt(18), WHITE)
    tx.text_frame.word_wrap = True
    
    # Placeholder for images
    slide.shapes.add_shape(1, Inches(6), Inches(1.5), Inches(3.5), Inches(2.5)).text = "Image Placeholder 1\n(Neon Lights)"
    slide.shapes.add_shape(1, Inches(6), Inches(4.5), Inches(3.5), Inches(2.5)).text = "Image Placeholder 2\n(Cyberpunk Vibe)"

    # 5. Concept 3
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, BLACK)
    add_text(slide, "CONCEPT 03: THE CRIMSON THEATRE", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(32), RED, True)
    
    details = "Mood: Dramatic, Cultural, Grand\n\n" \
              "Key Elements:\n" \
              "- Black Wood Structure (โครงไม้ทำสีดำ)\n" \
              "- Chinese Red Pillars (เสาสีแดงสด)\n" \
              "- Fabric Lanterns (โคมไฟผ้า)\n\n" \
              "References:\n" \
              "- ArchDaily: Jiu Li Hot Pot"
    
    tx = add_text(slide, details, Inches(0.5), Inches(1.5), Inches(5), Inches(5), Pt(18), WHITE)
    tx.text_frame.word_wrap = True
    
    # Placeholder for images
    slide.shapes.add_shape(1, Inches(6), Inches(1.5), Inches(3.5), Inches(2.5)).text = "Image Placeholder 1\n(Traditional Pillars)"
    slide.shapes.add_shape(1, Inches(6), Inches(4.5), Inches(3.5), Inches(2.5)).text = "Image Placeholder 2\n(Lanterns)"

    # 6. Conclusion
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, DARK_GREY)
    add_text(slide, "SUMMARY", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(36), RED, True)
    
    summary = "Concept 1: Modern Luxury -> ลูกค้า Premium / Couple\n" \
              "Concept 2: Neo-Cyber -> ลูกค้า Gen Z / Social Lovers\n" \
              "Concept 3: Crimson Theatre -> ลูกค้า Family / Mass"
              
    add_text(slide, summary, Inches(1), Inches(2.5), Inches(8), Inches(3), Pt(28), WHITE)

    prs.save('Mala_Hotpot_Moodboard.pptx')
    print("Presentation saved successfully.")

if __name__ == "__main__":
    create_presentation()
